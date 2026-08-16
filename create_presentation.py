import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ──────────────────────────────────────────────────────
# PATHS
# ──────────────────────────────────────────────────────
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
PLOT_DIR = os.path.join(BASE_DIR, "plots")
OUTPUT_PATH = os.path.join(BASE_DIR, "Master_Training_Project_Presentation.pptx")

IMG_STATES = os.path.join(PLOT_DIR, "closed_loop_states.png")
IMG_CONTROLS = os.path.join(PLOT_DIR, "control_actions.png")
IMG_ERROR = os.path.join(PLOT_DIR, "moving_horizon_error.png")

# ──────────────────────────────────────────────────────
# COLOR PALETTE
# ──────────────────────────────────────────────────────
NAVY       = RGBColor(12, 25, 58)
DARK_NAVY  = RGBColor(8, 16, 40)
BLUE       = RGBColor(30, 100, 200)
LIGHT_BLUE = RGBColor(70, 140, 230)
ACCENT     = RGBColor(0, 180, 216)       # Cyan accent
ACCENT2    = RGBColor(255, 170, 50)      # Amber accent
GREEN      = RGBColor(34, 180, 85)
RED_SOFT   = RGBColor(220, 70, 70)
WHITE      = RGBColor(255, 255, 255)
OFF_WHITE  = RGBColor(240, 243, 248)
LIGHT_BG   = RGBColor(245, 247, 252)
CARD_BG    = RGBColor(230, 237, 248)
DARK_GRAY  = RGBColor(55, 60, 72)
MID_GRAY   = RGBColor(120, 130, 150)
SUBTITLE   = RGBColor(170, 195, 230)

# ──────────────────────────────────────────────────────
# HELPERS
# ──────────────────────────────────────────────────────
def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_accent_bar(slide, x, y, w, h, color=ACCENT):
    """Add a thin colored accent bar."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar

def add_card(slide, x, y, w, h, fill_color=CARD_BG, border_color=None):
    """Add a rounded rectangle card."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    if border_color:
        card.line.color.rgb = border_color
        card.line.width = Pt(1.2)
    else:
        card.line.fill.background()
    return card

def set_tf_margins(tf, left=0.3, top=0.25, right=0.3, bottom=0.15):
    tf.margin_left = Inches(left)
    tf.margin_top = Inches(top)
    tf.margin_right = Inches(right)
    tf.margin_bottom = Inches(bottom)

def add_textbox(slide, x, y, w, h, text, size=14, color=DARK_GRAY, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tf

def add_bullet_card(slide, x, y, w, h, title, bullets, title_color=NAVY,
                    bullet_color=DARK_GRAY, card_color=CARD_BG, border=None,
                    title_size=18, bullet_size=13, icon="•"):
    """Add a card with title and bullet list."""
    card = add_card(slide, Inches(x), Inches(y), Inches(w), Inches(h), card_color, border)
    tf = card.text_frame
    tf.word_wrap = True
    set_tf_margins(tf)

    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(title_size)
    p.font.bold = True
    p.font.color.rgb = title_color
    p.space_after = Pt(12)

    for b in bullets:
        pb = tf.add_paragraph()
        pb.text = f"{icon}  {b}"
        pb.font.size = Pt(bullet_size)
        pb.font.color.rgb = bullet_color
        pb.space_after = Pt(8)
        pb.line_spacing = Pt(20)
    return card

def add_section_header(slide, title, subtitle="MASTER TRAINING PROJECT  •  IIT KHARAGPUR"):
    """Standard slide header with accent bar."""
    set_slide_bg(slide, WHITE)
    add_accent_bar(slide, Inches(0.8), Inches(0.55), Inches(0.28), Inches(0.9), ACCENT)

    add_textbox(slide, 1.3, 0.35, 11.0, 0.35, subtitle,
                size=10, color=BLUE, bold=True)
    add_textbox(slide, 1.3, 0.7, 11.0, 0.6, title,
                size=28, color=NAVY, bold=True)

    # Bottom rule
    add_accent_bar(slide, Inches(0.8), Inches(7.1), Inches(11.73), Pt(2), MID_GRAY)

def add_slide_number(slide, prs, num):
    """Add slide number bottom-right."""
    add_textbox(slide, 11.8, 7.05, 0.8, 0.35, str(num),
                size=10, color=MID_GRAY, align=PP_ALIGN.RIGHT)


# ──────────────────────────────────────────────────────
# BUILD PRESENTATION
# ──────────────────────────────────────────────────────
def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    slide_num = 0

    # ==================================================
    # SLIDE 1 — TITLE
    # ==================================================
    slide_num += 1
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, DARK_NAVY)

    # Decorative top accent line
    add_accent_bar(s, Inches(0), Inches(0), Inches(13.333), Pt(6), ACCENT)

    # Side accent strip
    add_accent_bar(s, Inches(0.8), Inches(1.4), Pt(5), Inches(4.5), ACCENT)

    # Title text
    tf = add_textbox(s, 1.3, 1.5, 10.5, 1.0,
                     "Physics-Informed Online Machine Learning",
                     size=38, color=WHITE, bold=True)
    p2 = tf.add_paragraph()
    p2.text = "for Model Predictive Control of Nonlinear Processes"
    p2.font.size = Pt(34)
    p2.font.bold = True
    p2.font.color.rgb = ACCENT
    p2.space_before = Pt(4)

    # Subtitle
    tf2 = add_textbox(s, 1.3, 3.2, 10.5, 0.5,
                      "Reproduction & Validation of PIRNN-LMPC Framework (Zheng & Wu, I&EC Research 2024)",
                      size=17, color=SUBTITLE)

    # Divider
    add_accent_bar(s, Inches(1.3), Inches(4.0), Inches(3.5), Pt(2), ACCENT)

    # Author info
    info_lines = [
        ("Devan Singh Faujdar", 20, WHITE, True),
        ("Indian Institute of Technology Kharagpur", 15, SUBTITLE, False),
        ("Department of Chemical Engineering", 13, MID_GRAY, False),
        ("August 2026", 13, MID_GRAY, False),
    ]
    y = 4.35
    for text, sz, clr, bld in info_lines:
        add_textbox(s, 1.3, y, 10.0, 0.4, text, size=sz, color=clr, bold=bld)
        y += 0.42

    # Right-side decorative badge
    badge = add_card(s, Inches(9.0), Inches(4.3), Inches(3.5), Inches(2.0), NAVY, ACCENT)
    btf = badge.text_frame
    btf.word_wrap = True
    set_tf_margins(btf, 0.3, 0.25)
    bp = btf.paragraphs[0]
    bp.text = "Target System"
    bp.font.size = Pt(11)
    bp.font.color.rgb = ACCENT
    bp.font.bold = True
    bp.space_after = Pt(6)
    bp2 = btf.add_paragraph()
    bp2.text = "Continuous Stirred-Tank\nReactor (CSTR)"
    bp2.font.size = Pt(17)
    bp2.font.color.rgb = WHITE
    bp2.font.bold = True
    bp2.space_after = Pt(6)
    bp3 = btf.add_paragraph()
    bp3.text = "with Parameter Drift & Uncertainty"
    bp3.font.size = Pt(12)
    bp3.font.color.rgb = SUBTITLE

    add_slide_number(s, prs, slide_num)

    # ==================================================
    # SLIDE 2 — PROBLEM STATEMENT
    # ==================================================
    slide_num += 1
    s = prs.slides.add_slide(blank)
    add_section_header(s, "Problem Statement & Motivation")

    add_bullet_card(s, 0.8, 1.6, 5.8, 5.3,
        "The Challenge",
        [
            "Data-driven surrogate models degrade when plant parameters drift from training conditions.",
            "Standard Neural Networks lack physical consistency — violate mass & energy conservation.",
            "Re-training from scratch is too slow for real-time MPC (needs < 36s sampling period).",
            "Need: A model that adapts online while respecting first-principles physics."
        ],
        card_color=RGBColor(255, 245, 235), border=ACCENT2, icon="⚠")

    add_bullet_card(s, 7.0, 1.6, 5.8, 5.3,
        "Our Approach: PIRNN-LMPC",
        [
            "Physics-Informed RNN embeds CSTR ODEs directly into the loss function as soft constraints.",
            "Error-triggered mechanism monitors prediction quality in real time.",
            "Online adaptation: jointly re-estimates model weights AND uncertain parameters (F, k₀).",
            "Lyapunov-based MPC guarantees closed-loop stability."
        ],
        card_color=RGBColor(230, 248, 240), border=GREEN, icon="✓")

    add_slide_number(s, prs, slide_num)

    # ==================================================
    # SLIDE 3 — CSTR PROCESS MODEL
    # ==================================================
    slide_num += 1
    s = prs.slides.add_slide(blank)
    add_section_header(s, "CSTR Process Model & Governing Equations")

    # ODE card
    ode_card = add_card(s, Inches(0.8), Inches(1.6), Inches(7.5), Inches(3.2), CARD_BG, BLUE)
    otf = ode_card.text_frame
    otf.word_wrap = True
    set_tf_margins(otf)

    op = otf.paragraphs[0]
    op.text = "Governing ODEs (Eq. 14a & 14b)"
    op.font.size = Pt(18)
    op.font.bold = True
    op.font.color.rgb = NAVY
    op.space_after = Pt(14)

    eqs = [
        "dCₐ/dt  =  (F / V) · (Cₐ₀ − Cₐ)  −  k₀ · exp(−E / RT) · Cₐ²",
        "",
        "dT/dt   =  (F / V) · (T₀ − T)  +  (−ΔH / ρCₚ) · k₀ · exp(−E / RT) · Cₐ²  +  Q / (ρCₚV)",
        "",
        "States:   x = [ΔCₐ, ΔT]    (deviations from steady state)",
        "Inputs:   u = [ΔCₐ₀, ΔQ]    (control deviations)"
    ]
    for eq in eqs:
        ep = otf.add_paragraph()
        ep.text = eq
        ep.font.size = Pt(13)
        ep.font.color.rgb = DARK_GRAY
        if eq.startswith("d"):
            ep.font.name = "Courier New"
            ep.font.bold = True
        ep.space_after = Pt(4)

    # Parameters table card
    param_card = add_card(s, Inches(8.7), Inches(1.6), Inches(3.83), Inches(3.2), NAVY)
    ptf = param_card.text_frame
    ptf.word_wrap = True
    set_tf_margins(ptf)

    pp = ptf.paragraphs[0]
    pp.text = "Process Parameters"
    pp.font.size = Pt(16)
    pp.font.bold = True
    pp.font.color.rgb = ACCENT
    pp.space_after = Pt(10)

    params = [
        ("F₀", "5.0 m³/h"),
        ("V", "1.0 m³"),
        ("k₀", "8.46 × 10⁶"),
        ("E", "5.0 × 10⁴ kJ/kmol"),
        ("T₀", "300 K"),
        ("ΔH", "−1.15 × 10⁴ kJ/kmol"),
        ("ρ_L", "1000 kg/m³"),
        ("Cₚ", "0.231 kJ/(kg·K)"),
    ]
    for sym, val in params:
        lp = ptf.add_paragraph()
        lp.text = f"{sym:>6}  =  {val}"
        lp.font.size = Pt(11)
        lp.font.color.rgb = WHITE
        lp.font.name = "Courier New"
        lp.space_after = Pt(3)

    # Steady state card
    ss_card = add_card(s, Inches(0.8), Inches(5.1), Inches(5.8), Inches(1.8), RGBColor(230, 248, 240), GREEN)
    stf = ss_card.text_frame
    stf.word_wrap = True
    set_tf_margins(stf)
    sp = stf.paragraphs[0]
    sp.text = "Steady-State Operating Point"
    sp.font.size = Pt(16)
    sp.font.bold = True
    sp.font.color.rgb = NAVY
    sp.space_after = Pt(8)

    ss_items = [
        "Cₐₛ = 1.95 kmol/m³    |    Tₛ = 402.0 K    |    Cₐ₀ₛ = 4.0 kmol/m³    |    Qₛ = 0.0 kJ/h",
        "Control Bounds:  |ΔCₐ₀| ≤ 3.5 kmol/m³,    |ΔQ| ≤ 5 × 10⁵ kJ/h"
    ]
    for item in ss_items:
        ip = stf.add_paragraph()
        ip.text = item
        ip.font.size = Pt(12)
        ip.font.color.rgb = DARK_GRAY
        ip.space_after = Pt(4)

    # Disturbance schedule card
    dist_card = add_card(s, Inches(7.0), Inches(5.1), Inches(5.53), Inches(1.8), RGBColor(255, 240, 240), RED_SOFT)
    dtf = dist_card.text_frame
    dtf.word_wrap = True
    set_tf_margins(dtf)
    dp = dtf.paragraphs[0]
    dp.text = "Disturbance Schedule"
    dp.font.size = Pt(16)
    dp.font.bold = True
    dp.font.color.rgb = NAVY
    dp.space_after = Pt(8)

    dists = [
        "t = 0.09 h:  F → 8.0 m³/h (160%)    k₀ → 6.77×10⁶ (80%)    [Moderate]",
        "t = 0.19 h:  F → 11.5 m³/h (230%)   k₀ → 2.54×10⁶ (30%)    [Severe]"
    ]
    for d in dists:
        ddp = dtf.add_paragraph()
        ddp.text = d
        ddp.font.size = Pt(12)
        ddp.font.color.rgb = DARK_GRAY
        ddp.space_after = Pt(4)

    add_slide_number(s, prs, slide_num)

    # ==================================================
    # SLIDE 4 — PIRNN ARCHITECTURE
    # ==================================================
    slide_num += 1
    s = prs.slides.add_slide(blank)
    add_section_header(s, "PIRNN Model Architecture")

    # Architecture flow card
    arch_card = add_card(s, Inches(0.8), Inches(1.6), Inches(8.0), Inches(2.8), CARD_BG, BLUE)
    atf = arch_card.text_frame
    atf.word_wrap = True
    set_tf_margins(atf)

    ap = atf.paragraphs[0]
    ap.text = "Data Flow:  Input → Encoder → Recurrent Core → Decoder → Output"
    ap.font.size = Pt(16)
    ap.font.bold = True
    ap.font.color.rgb = NAVY
    ap.space_after = Pt(14)

    arch_items = [
        ("Input", "[x₀, u] ∈ ℝ⁴  →  initial state deviation (2) + control deviation (2)"),
        ("Encoder", "Linear(4→64) → Tanh → Linear(64→64) → Tanh"),
        ("GRU Cell", "GRUCell(64, 64) applied 10 times  —  self-recurrent: h = GRU(h, h)"),
        ("Decoder", "Linear(64→32) → Tanh → Linear(32→2)"),
        ("Residual", "xₜ₊₁ = xₜ + δ(hₜ)   — each sub-step adds learned increment"),
        ("Output", "Trajectory tensor: (batch, 11, 2)  — 10 sub-steps + initial state"),
    ]
    for label, desc in arch_items:
        lp = atf.add_paragraph()
        lp.text = f"  {label:>12}  │  {desc}"
        lp.font.size = Pt(12)
        lp.font.color.rgb = DARK_GRAY
        lp.font.name = "Courier New"
        lp.space_after = Pt(3)

    # Parameter count card
    pc_card = add_card(s, Inches(9.2), Inches(1.6), Inches(3.33), Inches(2.8), NAVY)
    pctf = pc_card.text_frame
    pctf.word_wrap = True
    set_tf_margins(pctf)
    pcp = pctf.paragraphs[0]
    pcp.text = "Parameter Count"
    pcp.font.size = Pt(16)
    pcp.font.bold = True
    pcp.font.color.rgb = ACCENT
    pcp.space_after = Pt(10)

    layers = [
        ("Encoder", "4,544"),
        ("GRU Cell", "24,960"),
        ("Decoder", "2,082"),
        ("─────────", "──────"),
        ("TOTAL", "31,586"),
    ]
    for name, count in layers:
        lcp = pctf.add_paragraph()
        lcp.text = f"{name:<12} {count:>8}"
        lcp.font.size = Pt(12)
        lcp.font.color.rgb = WHITE
        lcp.font.name = "Courier New"
        lcp.space_after = Pt(3)
        if name == "TOTAL":
            lcp.font.bold = True
            lcp.font.color.rgb = ACCENT

    # Loss function card
    loss_card = add_card(s, Inches(0.8), Inches(4.7), Inches(5.8), Inches(2.2), RGBColor(255, 245, 235), ACCENT2)
    ltf = loss_card.text_frame
    ltf.word_wrap = True
    set_tf_margins(ltf)
    lp = ltf.paragraphs[0]
    lp.text = "Hybrid Loss Function (Eq. 5 & 7)"
    lp.font.size = Pt(16)
    lp.font.bold = True
    lp.font.color.rgb = NAVY
    lp.space_after = Pt(10)

    loss_items = [
        "L_total  =  L_data  +  η · L_physics",
        "",
        "L_data:    MSE between predicted & RK4-generated trajectories",
        "L_physics: Residual of CSTR ODEs via finite differences",
        "η = 0.1    (physics loss weight coefficient)",
        "Temp residuals scaled by 1/100 for magnitude balance"
    ]
    for li in loss_items:
        lip = ltf.add_paragraph()
        lip.text = li
        lip.font.size = Pt(12)
        lip.font.color.rgb = DARK_GRAY
        if li.startswith("L_total"):
            lip.font.size = Pt(14)
            lip.font.bold = True
            lip.font.name = "Courier New"
        lip.space_after = Pt(3)

    # Training config card
    train_card = add_card(s, Inches(7.0), Inches(4.7), Inches(5.53), Inches(2.2), CARD_BG, BLUE)
    ttf = train_card.text_frame
    ttf.word_wrap = True
    set_tf_margins(ttf)
    tp = ttf.paragraphs[0]
    tp.text = "Training Configuration"
    tp.font.size = Pt(16)
    tp.font.bold = True
    tp.font.color.rgb = NAVY
    tp.space_after = Pt(10)

    train_items = [
        ("Collocation Data", "80,000 trajectories (200 states × 400 inputs)"),
        ("Integration", "Vectorized PyTorch RK4 — generated in < 0.5s"),
        ("Split", "60% train / 20% val / 20% test"),
        ("Optimizer", "Adam, lr = 1×10⁻³"),
        ("Batch Size", "2,048"),
        ("Epochs", "10 (early stopping, patience = 5)"),
        ("Device", "Apple MPS (Metal GPU)"),
    ]
    for label, val in train_items:
        tip = ttf.add_paragraph()
        tip.text = f"  {label}: {val}"
        tip.font.size = Pt(11)
        tip.font.color.rgb = DARK_GRAY
        tip.space_after = Pt(3)

    add_slide_number(s, prs, slide_num)

    # ==================================================
    # SLIDE 5 — ERROR-TRIGGERED MECHANISM & MPC
    # ==================================================
    slide_num += 1
    s = prs.slides.add_slide(blank)
    add_section_header(s, "Error-Triggered Online Update & MPC Controller")

    add_bullet_card(s, 0.8, 1.6, 5.8, 2.8,
        "Error-Triggered Mechanism",
        [
            "Prediction error: E_RNN(t) = MSE(PIRNN prediction, plant state)",
            "Threshold: E_T = 1 × 10⁻⁴",
            "Trigger: E_RNN > E_T for 3 consecutive steps → online update",
            "Update: 50 Adam steps (lr = 10⁻³) on recent history (5 states, 4 inputs)",
            "Joint estimation: model weights AND process parameters (F, k₀)"
        ],
        icon="→")

    add_bullet_card(s, 7.0, 1.6, 5.53, 2.8,
        "MPC Objective Function",
        [
            "min  w_CA·ΔCₐ²  +  w_T·(ΔT/100)²  +  w_u1·ΔCₐ₀²  +  w_u2·ΔQ²",
            "w_CA = 10.0  |  w_T = 1.0  |  w_u1 = 0.1  |  w_u2 = 10⁻¹⁰",
            "Solver: SLSQP, maxiter = 50, ftol = 10⁻⁴",
            "Single-step prediction horizon"
        ],
        icon="▸")

    # Three comparison schemes
    schemes = [
        ("PIRNN_no_update", "Static Baseline", "No online updates.\nNominal model used\nthroughout.", RED_SOFT),
        ("PIRNN_data_online", "Data-Only Online", "Error-triggered update\nwith data MSE loss\nonly (no physics).", ACCENT2),
        ("PIRNN_physics_enhanced", "Proposed Method", "Error-triggered update\nwith hybrid loss:\nL_data + η·L_physics.", GREEN),
    ]

    for i, (scheme, label, desc, color) in enumerate(schemes):
        x = 0.8 + i * 4.1
        card = add_card(s, Inches(x), Inches(4.8), Inches(3.7), Inches(2.1), WHITE, color)
        tf = card.text_frame
        tf.word_wrap = True
        set_tf_margins(tf)

        p1 = tf.paragraphs[0]
        p1.text = f"Scheme {i+1}"
        p1.font.size = Pt(10)
        p1.font.color.rgb = color
        p1.font.bold = True
        p1.space_after = Pt(4)

        p2 = tf.add_paragraph()
        p2.text = label
        p2.font.size = Pt(16)
        p2.font.bold = True
        p2.font.color.rgb = NAVY
        p2.space_after = Pt(8)

        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(12)
        p3.font.color.rgb = DARK_GRAY

    add_slide_number(s, prs, slide_num)

    # ==================================================
    # SLIDE 6 — FIGURE 4: STATE TRAJECTORIES
    # ==================================================
    slide_num += 1
    s = prs.slides.add_slide(blank)
    add_section_header(s, "Results — Closed-Loop State Trajectories (Figure 4)")

    if os.path.exists(IMG_STATES):
        s.shapes.add_picture(IMG_STATES, Inches(0.6), Inches(1.5), width=Inches(8.5))

    obs_card = add_card(s, Inches(9.4), Inches(1.5), Inches(3.4), Inches(5.5), CARD_BG, BLUE)
    otf = obs_card.text_frame
    otf.word_wrap = True
    set_tf_margins(otf)

    op = otf.paragraphs[0]
    op.text = "Key Observations"
    op.font.size = Pt(16)
    op.font.bold = True
    op.font.color.rgb = NAVY
    op.space_after = Pt(10)

    obs = [
        "Pre-disturbance: All 3 schemes hold near-zero deviation — nominal PIRNN is accurate.",
        "Post t=0.09h: All schemes begin deviating as plant parameters shift.",
        "Post t=0.19h: Severe disturbance (F=230%, k₀=30%) causes large state excursions.",
        "Online schemes achieve better ΔT control (6.5% IAE improvement) but larger ΔCₐ deviation.",
        "Physics-enhanced and data-only schemes are nearly identical visually."
    ]
    for o in obs:
        op2 = otf.add_paragraph()
        op2.text = f"▸  {o}"
        op2.font.size = Pt(11)
        op2.font.color.rgb = DARK_GRAY
        op2.space_after = Pt(8)

    add_slide_number(s, prs, slide_num)

    # ==================================================
    # SLIDE 7 — FIGURE 5: CONTROL ACTIONS
    # ==================================================
    slide_num += 1
    s = prs.slides.add_slide(blank)
    add_section_header(s, "Results — Control Actions Applied (Figure 5)")

    if os.path.exists(IMG_CONTROLS):
        s.shapes.add_picture(IMG_CONTROLS, Inches(0.6), Inches(1.5), width=Inches(8.5))

    obs_card = add_card(s, Inches(9.4), Inches(1.5), Inches(3.4), Inches(5.5), CARD_BG, BLUE)
    otf = obs_card.text_frame
    otf.word_wrap = True
    set_tf_margins(otf)

    op = otf.paragraphs[0]
    op.text = "Key Observations"
    op.font.size = Pt(16)
    op.font.bold = True
    op.font.color.rgb = NAVY
    op.space_after = Pt(10)

    obs = [
        "All control inputs remain within physical saturation bounds.",
        "Control magnitudes are small: ΔCₐ₀ ≈ 10⁻³, ΔQ ≈ 10⁻⁶ × 10⁵.",
        "No-update scheme produces most oscillatory feed concentration control.",
        "Largest ΔQ spike from static model after second disturbance.",
        "Online schemes produce smoother, more conservative actions."
    ]
    for o in obs:
        op2 = otf.add_paragraph()
        op2.text = f"▸  {o}"
        op2.font.size = Pt(11)
        op2.font.color.rgb = DARK_GRAY
        op2.space_after = Pt(8)

    add_slide_number(s, prs, slide_num)

    # ==================================================
    # SLIDE 8 — FIGURE 6: PREDICTION ERROR
    # ==================================================
    slide_num += 1
    s = prs.slides.add_slide(blank)
    add_section_header(s, "Results — Moving-Horizon Prediction Error (Figure 6)")

    if os.path.exists(IMG_ERROR):
        s.shapes.add_picture(IMG_ERROR, Inches(0.6), Inches(1.5), width=Inches(8.5))

    obs_card = add_card(s, Inches(9.4), Inches(1.5), Inches(3.4), Inches(5.5), CARD_BG, BLUE)
    otf = obs_card.text_frame
    otf.word_wrap = True
    set_tf_margins(otf)

    op = otf.paragraphs[0]
    op.text = "Key Observations"
    op.font.size = Pt(16)
    op.font.bold = True
    op.font.color.rgb = NAVY
    op.space_after = Pt(10)

    obs = [
        "All errors > threshold E_T = 10⁻⁴ — error-triggered mechanism fires frequently.",
        "Static model error peaks at ~30 MSE — confirms model-plant mismatch.",
        "Data-online achieves lowest mean error: 1.22 (vs. 18.97 static).",
        "Physics-enhanced mean error: 1.55 — 92% reduction vs. static baseline.",
        "Sharp spikes at disturbance onset times with partial recovery for online schemes."
    ]
    for o in obs:
        op2 = otf.add_paragraph()
        op2.text = f"▸  {o}"
        op2.font.size = Pt(11)
        op2.font.color.rgb = DARK_GRAY
        op2.space_after = Pt(8)

    add_slide_number(s, prs, slide_num)

    # ==================================================
    # SLIDE 9 — QUANTITATIVE METRICS TABLE
    # ==================================================
    slide_num += 1
    s = prs.slides.add_slide(blank)
    add_section_header(s, "Quantitative Performance Metrics")

    # Main metrics table
    table_data = [
        ["Metric", "PIRNN_no_update", "PIRNN_data_online", "PIRNN_physics_enhanced"],
        ["IAE (Cₐ)", "0.0811  ★", "0.1896", "0.1914"],
        ["IAE (T)", "10.1515", "9.4938", "9.4886  ★"],
        ["ISE (Cₐ)", "0.0571  ★", "0.2246", "0.2294"],
        ["ISE (T)", "630.79", "562.94", "562.36  ★"],
        ["Mean E_RNN", "18.97", "1.22  ★", "1.55"],
    ]

    tbl = s.shapes.add_table(len(table_data), 4,
                              Inches(0.8), Inches(1.6),
                              Inches(11.73), Inches(3.2)).table

    # Style header
    for j in range(4):
        cell = tbl.cell(0, j)
        cell.text = table_data[0][j]
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY

    # Style body
    for i in range(1, len(table_data)):
        for j in range(4):
            cell = tbl.cell(i, j)
            cell.text = table_data[i][j]
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = DARK_GRAY
                p.alignment = PP_ALIGN.CENTER
                if "★" in table_data[i][j]:
                    p.font.bold = True
                    p.font.color.rgb = GREEN
            cell.fill.solid()
            cell.fill.fore_color.rgb = OFF_WHITE if i % 2 == 0 else WHITE

    # Improvement card
    imp_card = add_card(s, Inches(0.8), Inches(5.1), Inches(5.8), Inches(1.8), RGBColor(230, 248, 240), GREEN)
    itf = imp_card.text_frame
    itf.word_wrap = True
    set_tf_margins(itf)
    ip = itf.paragraphs[0]
    ip.text = "✓  Physics-Enhanced vs. Static Baseline"
    ip.font.size = Pt(16)
    ip.font.bold = True
    ip.font.color.rgb = NAVY
    ip.space_after = Pt(8)

    improvements = [
        "Temperature IAE:  +6.53% improvement  (better tracking)",
        "Temperature ISE:  +10.85% improvement  (lower peak errors)",
        "Prediction Error:  +91.83% reduction  (model accuracy)",
    ]
    for imp in improvements:
        iip = itf.add_paragraph()
        iip.text = f"  ✓  {imp}"
        iip.font.size = Pt(12)
        iip.font.color.rgb = DARK_GRAY
        iip.space_after = Pt(4)

    # Trade-off card
    trade_card = add_card(s, Inches(7.0), Inches(5.1), Inches(5.53), Inches(1.8), RGBColor(255, 240, 240), RED_SOFT)
    trtf = trade_card.text_frame
    trtf.word_wrap = True
    set_tf_margins(trtf)
    trp = trtf.paragraphs[0]
    trp.text = "⚠  Concentration Control Trade-off"
    trp.font.size = Pt(16)
    trp.font.bold = True
    trp.font.color.rgb = NAVY
    trp.space_after = Pt(8)

    trade_items = [
        "Concentration IAE:  −136% (larger deviation than static)",
        "Online updates alter model landscape, affecting MPC trajectory",
        "Root cause: MPC weight tuning & under-utilized control authority"
    ]
    for ti in trade_items:
        tip = trtf.add_paragraph()
        tip.text = f"  ⚠  {ti}"
        tip.font.size = Pt(12)
        tip.font.color.rgb = DARK_GRAY
        tip.space_after = Pt(4)

    add_slide_number(s, prs, slide_num)

    # ==================================================
    # SLIDE 10 — PARAMETER ESTIMATION
    # ==================================================
    slide_num += 1
    s = prs.slides.add_slide(blank)
    add_section_header(s, "Online Parameter Estimation Accuracy")

    # Estimation table
    est_data = [
        ["Phase", "True F (m³/h)", "Est. F", "True k₀", "Est. k₀", "Accuracy"],
        ["Nominal (t < 0.09h)", "5.00", "5.00", "8.46×10⁶", "8.46×10⁶", "100.00%"],
        ["Moderate (t = 0.09h)", "8.00", "7.99–8.00", "6.77×10⁶", "6.77×10⁶", ">99.88%"],
        ["Severe (t = 0.19h)", "11.50", "11.49–11.50", "2.54×10⁶", "2.54×10⁶", ">99.91%"],
    ]

    tbl = s.shapes.add_table(len(est_data), 6,
                              Inches(0.8), Inches(1.6),
                              Inches(11.73), Inches(2.2)).table

    for j in range(6):
        cell = tbl.cell(0, j)
        cell.text = est_data[0][j]
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY

    for i in range(1, len(est_data)):
        for j in range(6):
            cell = tbl.cell(i, j)
            cell.text = est_data[i][j]
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = DARK_GRAY
                p.alignment = PP_ALIGN.CENTER
                if j == 5:  # accuracy column
                    p.font.bold = True
                    p.font.color.rgb = GREEN
            cell.fill.solid()
            cell.fill.fore_color.rgb = OFF_WHITE if i % 2 == 0 else WHITE

    # Timing card
    add_bullet_card(s, 0.8, 4.2, 5.8, 2.7,
        "⏱  Online Update Timing",
        [
            "Data-only updates:  ~1.17s mean  (10 total updates)",
            "Physics-enhanced:  ~1.22s mean  (10 total updates)",
            "Physics overhead:  ~4% additional compute",
            "Sampling period limit:  36s  →  ample real-time margin",
            "Both schemes:  total online time < 13 seconds"
        ],
        card_color=CARD_BG, border=BLUE, icon="▸", bullet_size=12)

    # Key takeaway
    key_card = add_card(s, Inches(7.0), Inches(4.2), Inches(5.53), Inches(2.7), RGBColor(230, 248, 240), GREEN)
    ktf = key_card.text_frame
    ktf.word_wrap = True
    set_tf_margins(ktf)

    kp = ktf.paragraphs[0]
    kp.text = "🏆  Key Result"
    kp.font.size = Pt(18)
    kp.font.bold = True
    kp.font.color.rgb = NAVY
    kp.space_after = Pt(12)

    key_texts = [
        "Both online schemes achieve >99.8% parameter estimation accuracy within 50 gradient steps (~1s).",
        "",
        "Even under severe disturbance (F at 230%, k₀ at 30% of nominal), the joint optimization recovers true parameters to 4 significant figures.",
        "",
        "This validates the PIRNN as a viable real-time surrogate for online MPC."
    ]
    for kt in key_texts:
        ktp = ktf.add_paragraph()
        ktp.text = kt
        ktp.font.size = Pt(12)
        ktp.font.color.rgb = DARK_GRAY
        ktp.space_after = Pt(4)

    add_slide_number(s, prs, slide_num)

    # ==================================================
    # SLIDE 11 — ANALYSIS & DISCUSSION
    # ==================================================
    slide_num += 1
    s = prs.slides.add_slide(blank)
    add_section_header(s, "Analysis & Discussion")

    add_bullet_card(s, 0.8, 1.6, 5.8, 2.6,
        "Concentration vs. Temperature Trade-off",
        [
            "Online update modifies model to prioritize thermal dynamics (coupled through Arrhenius term).",
            "MPC weights (w_CA=10, w_T=1) penalize concentration heavily, but updated model's altered prediction landscape shifts controller behavior.",
            "Result: better temperature tracking at the expense of concentration."
        ],
        icon="▸", card_color=RGBColor(255, 245, 235), border=ACCENT2)

    add_bullet_card(s, 7.0, 1.6, 5.53, 2.6,
        "Physics vs. Data-Only: Why So Similar?",
        [
            "IAE difference: < 1% between the two online schemes.",
            "With only 50 gradient steps and 5-state history window, physics regularization has limited opportunity to differentiate.",
            "Parameter estimation accuracy is equally excellent for both — the data signal dominates."
        ],
        icon="▸", card_color=CARD_BG, border=BLUE)

    add_bullet_card(s, 0.8, 4.5, 5.8, 2.5,
        "Prediction Error Dynamics",
        [
            "All errors remain >> threshold E_T = 10⁻⁴ → continuous triggering.",
            "92% mean error reduction with online updates validates the mechanism.",
            "Model capacity (31K params) may be insufficient for severe disturbance regime.",
        ],
        icon="▸", card_color=CARD_BG, border=BLUE)

    add_bullet_card(s, 7.0, 4.5, 5.53, 2.5,
        "Control Authority Under-Utilization",
        [
            "Control actions are 3–4 orders of magnitude below physical limits.",
            "ΔCₐ₀ ≈ 10⁻³ vs. limit of 3.5  |  ΔQ ≈ 10⁻¹ vs. limit of 5×10⁵.",
            "SLSQP with 50 iterations may converge to local minima near zero initial guess."
        ],
        icon="⚠", card_color=RGBColor(255, 240, 240), border=RED_SOFT)

    add_slide_number(s, prs, slide_num)

    # ==================================================
    # SLIDE 12 — CODEBASE OVERVIEW
    # ==================================================
    slide_num += 1
    s = prs.slides.add_slide(blank)
    add_section_header(s, "Implementation & Codebase Overview")

    files = [
        ("config.py", "57 lines", "All hyperparameters & CSTR physical constants"),
        ("cstr_env.py", "88 lines", "Plant simulator — SciPy ODE integration"),
        ("pirnn_model.py", "104 lines", "PIRNN architecture + physics residual + hybrid loss"),
        ("train_pirnn.py", "176 lines", "Offline training with vectorized RK4 data generation"),
        ("mpc_controller.py", "65 lines", "PIRNN-based MPC with SLSQP optimization"),
        ("run_closed_loop_simulation.py", "236 lines", "3-scheme comparison simulation + plotting"),
        ("evaluate_metrics.py", "70 lines", "Quantitative IAE/ISE metric computation"),
    ]

    # File table
    tbl = s.shapes.add_table(len(files) + 1, 3,
                              Inches(0.8), Inches(1.6),
                              Inches(7.5), Inches(3.8)).table

    headers = ["File", "Size", "Description"]
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY

    for i, (fname, size, desc) in enumerate(files):
        for j, val in enumerate([fname, size, desc]):
            cell = tbl.cell(i + 1, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = DARK_GRAY
                p.alignment = PP_ALIGN.LEFT if j == 2 else PP_ALIGN.CENTER
                if j == 0:
                    p.font.name = "Courier New"
                    p.font.bold = True
            cell.fill.solid()
            cell.fill.fore_color.rgb = OFF_WHITE if i % 2 == 1 else WHITE

    # Tech stack card
    add_bullet_card(s, 8.7, 1.6, 3.83, 3.8,
        "Tech Stack",
        [
            "Python 3.x",
            "PyTorch (MPS backend)",
            "SciPy (solve_ivp, minimize)",
            "NumPy",
            "Matplotlib",
            "python-pptx",
            "",
            "Total:  ~796 lines",
            "Model:  128.1 KB checkpoint",
        ],
        card_color=NAVY, title_color=ACCENT, bullet_color=WHITE, icon="•", bullet_size=12)

    # Stats badges
    stats = [
        ("31,586", "Model Parameters"),
        ("80,000", "Training Trajectories"),
        ("128 KB", "Checkpoint Size"),
        ("~1s", "Online Update Speed"),
    ]

    for i, (num, label) in enumerate(stats):
        x = 0.8 + i * 3.1
        badge = add_card(s, Inches(x), Inches(5.7), Inches(2.7), Inches(1.2), CARD_BG, BLUE)
        btf = badge.text_frame
        btf.word_wrap = True
        set_tf_margins(btf, 0.2, 0.15)
        bp1 = btf.paragraphs[0]
        bp1.text = num
        bp1.font.size = Pt(24)
        bp1.font.bold = True
        bp1.font.color.rgb = BLUE
        bp1.alignment = PP_ALIGN.CENTER
        bp1.space_after = Pt(2)
        bp2 = btf.add_paragraph()
        bp2.text = label
        bp2.font.size = Pt(11)
        bp2.font.color.rgb = DARK_GRAY
        bp2.alignment = PP_ALIGN.CENTER

    add_slide_number(s, prs, slide_num)

    # ==================================================
    # SLIDE 13 — RECOMMENDATIONS
    # ==================================================
    slide_num += 1
    s = prs.slides.add_slide(blank)
    add_section_header(s, "Recommendations & Future Work")

    add_bullet_card(s, 0.8, 1.6, 5.8, 5.3,
        "Immediate Improvements",
        [
            "Multi-step MPC horizon — use N-step prediction instead of single-step for better foresight.",
            "Rebalance MPC weights — w_u2 = 10⁻¹⁰ effectively removes heat input penalty; should be tuned.",
            "Expand online history window — 5 states is minimal; 15–20 would give richer adaptation signal.",
            "Increase model capacity — hidden_size = 128 or 2-layer GRU for severe disturbance regimes.",
            "Adaptive error threshold — replace fixed E_T = 10⁻⁴ with moving-average baseline to prevent continuous triggering.",
            "More training epochs — 10 epochs may be insufficient for full convergence; try 50–100 with scheduler."
        ],
        icon="→", card_color=CARD_BG, border=BLUE)

    add_bullet_card(s, 7.0, 1.6, 5.53, 5.3,
        "Master Thesis Extensions",
        [
            "Multi-reactor networks — extend PIRNN-MPC to cascade reactor systems and distillation columns.",
            "Measurement noise — integrate Extended Kalman Filter (EKF) for noisy sensor measurements.",
            "Continuous PINNs — test space-time PINNs with automatic differentiation vs. finite-difference residuals.",
            "Hardware-in-the-Loop — deploy controller on embedded edge devices for real-time HIL validation.",
            "Robust MPC — add stochastic uncertainty sets around estimated parameters for safety guarantees.",
            "Transfer learning — pre-train on multiple reactor configurations for faster online adaptation."
        ],
        icon="🚀", card_color=RGBColor(230, 248, 240), border=GREEN)

    add_slide_number(s, prs, slide_num)

    # ==================================================
    # SLIDE 14 — THANK YOU / END
    # ==================================================
    slide_num += 1
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, DARK_NAVY)

    add_accent_bar(s, Inches(0), Inches(0), Inches(13.333), Pt(6), ACCENT)
    add_accent_bar(s, Inches(0.8), Inches(2.2), Pt(5), Inches(3.0), ACCENT)

    add_textbox(s, 1.3, 2.4, 10.5, 0.8,
                "Thank You", size=44, color=WHITE, bold=True)

    tf = add_textbox(s, 1.3, 3.4, 10.5, 0.5,
                     "Physics-Informed Online ML for Model Predictive Control",
                     size=20, color=ACCENT)

    add_accent_bar(s, Inches(1.3), Inches(4.2), Inches(3.0), Pt(2), ACCENT)

    info = [
        ("Devan Singh Faujdar", 18, WHITE, True),
        ("IIT Kharagpur  •  Department of Chemical Engineering", 14, SUBTITLE, False),
        ("", 10, SUBTITLE, False),
        ("Questions & Discussion Welcome", 16, ACCENT, True),
    ]
    y = 4.5
    for text, sz, clr, bld in info:
        add_textbox(s, 1.3, y, 10.0, 0.4, text, size=sz, color=clr, bold=bld)
        y += 0.45

    # Summary badge
    badge = add_card(s, Inches(8.5), Inches(2.5), Inches(4.0), Inches(3.5), NAVY, ACCENT)
    btf = badge.text_frame
    btf.word_wrap = True
    set_tf_margins(btf)

    bp = btf.paragraphs[0]
    bp.text = "Key Achievements"
    bp.font.size = Pt(14)
    bp.font.bold = True
    bp.font.color.rgb = ACCENT
    bp.space_after = Pt(10)

    achievements = [
        "✓  Full PIRNN-LMPC framework implemented",
        "✓  >99.8% parameter estimation accuracy",
        "✓  92% prediction error reduction",
        "✓  6.5% temperature control improvement",
        "✓  Real-time feasible (~1s updates)",
        "✓  3 comparison schemes evaluated",
        "✓  Figures 4, 5, 6 reproduced",
    ]
    for a in achievements:
        ap = btf.add_paragraph()
        ap.text = a
        ap.font.size = Pt(12)
        ap.font.color.rgb = WHITE
        ap.space_after = Pt(5)

    add_slide_number(s, prs, slide_num)

    # ──────────────────────────────────────────────
    # SAVE
    # ──────────────────────────────────────────────
    prs.save(OUTPUT_PATH)
    print(f"\n{'='*60}")
    print(f"  Presentation saved successfully!")
    print(f"  Path: {OUTPUT_PATH}")
    print(f"  Slides: {slide_num}")
    print(f"{'='*60}")


if __name__ == "__main__":
    build_presentation()
